import dash
from dash import html, dcc, Input, Output
import dash_bootstrap_components as dbc
import plotly.express as px
import pandas as pd
import plotly.graph_objects as go
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a Dash web application
app = dash.Dash(__name__)
data = px.data.iris()
fig = px.scatter(data, x='sepal_width', y='sepal_length',
                 color='species', title='Scatter Plot')

fig2 = px.bar(data, x='sepal_width', y='sepal_length',
              color='species', title='bar Plot')

nodes = ["Node A", "Node B", "Node C", "Node D"]
links = [
    {"source": 0, "target": 1, "value": 20},
    {"source": 1, "target": 2, "value": 10},
    {"source": 1, "target": 3, "value": 10},
]
aspect_ratio = len(nodes) / (len(links) + 1)
# Define the layout of the app
app.layout = html.Div(
    children=[
        html.H1("test App"),
        dbc.Row([
            dbc.Col([
                html.Img(src='assets/template1.png', alt='Image'),
                dbc.Button("Choose this template", id="template"),
                html.Div(id='notice', children=[]),
                dcc.Graph(
                    id='example-graph',
                    figure=fig
                ),
                dcc.Graph(
                    id='example-graph2',
                    figure=fig2
                ),
                dcc.Graph(
                    id='sankey-chart',
                    figure={
                        'data': [go.Sankey(
                            node=dict(
                                pad=15,
                                thickness=20,
                                line=dict(color='black', width=0.5),
                                label=nodes
                            ),
                            link=dict(
                                source=[link['source'] for link in links],
                                target=[link['target'] for link in links],
                                value=[link['value'] for link in links]
                            )
                        )],
                        'layout': go.Layout(
                            title='Square Sankey Chart',
                            # aspectratio=dict(x=1, y=aspect_ratio),
                            autosize=False
                        )
                    }
                ),
                dbc.Button('Add to pptx', id='add_btn'),
                html.Div(id='output')
            ])
        ])
    ]
)


@app.callback(
    Output('output', 'children'),
    Input('add_btn', 'n_clicks'),
)
def toggle(n_clicks):
    filename = "test2.pptx"

    pr = Presentation("template1.pptx")

    slides = pr.slides
    slides._sldIdLst.clear()
    slide_layout = pr.slide_layouts[5]
    slide = pr.slides.add_slide(slide_layout)
    image_path = "newplot.png"
    left = Inches(.1)
    top = Inches(.1)
    # width = Inches(13)
    height = Inches(7)
    # max height and width
    slide.shapes.add_picture(image_path, left, top, None, height)

    pr.save(filename)
    return html.Div("printed")


if __name__ == '__main__':
    app.run_server(debug=True)
